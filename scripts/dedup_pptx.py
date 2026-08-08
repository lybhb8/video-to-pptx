#!/usr/bin/env python3
"""
dedup_pptx.py
Merge duplicate slides in an existing .pptx deck without touching the source
video. Two signals, either one of which drops a slide (compared against the
previously kept slide):

1. Image: perceptual hash of the slide's embedded frame image, optionally
   cropped to the main content region (--crop-*) to ignore overlay strips
   such as title bars, side panels, taskbars or watermarks.
2. Text: OCR the slide image, keep only text inside the main content region
   (--region-top/bottom/right, --min-text-height, --min-conf), and drop the
   slide when most of the text is the same (--text-similarity, e.g. 0.50).

Usage:
  python scripts/dedup_pptx.py input.pptx output.pptx \
      --image-similarity 0.95 --text-similarity 0.5 \
      --crop-top 0.15 --crop-bottom 0.08 --crop-right 0.10 \
      --region-top 0.15 --region-bottom 0.92 --region-right 0.90
"""
import argparse
import hashlib
import sys
from io import BytesIO
from pathlib import Path

import numpy as np
from PIL import Image
from pptx import Presentation

from video_to_pptx import hash_similarity, text_similarity


def image_hash_pil(img: Image.Image, size: int = 16, crop=None) -> str:
    """Perceptual hash of a PIL image; `crop` = (top, bottom, left, right)
    fractions trimmed from each edge before hashing (overlay removal)."""
    img = img.convert("L")
    if crop:
        w, h = img.size
        top, bottom, left, right = crop
        img = img.crop((int(w * left), int(h * top),
                        int(w * (1 - right)), int(h * (1 - bottom))))
    img = img.resize((size, size), Image.Resampling.LANCZOS)
    arr = np.array(img, dtype=np.float32)
    bits = (arr > arr.mean()).astype(np.uint8)
    return hashlib.sha1(bits.tobytes()).hexdigest()


_ocr_reader = None


def ocr_content_text(img: Image.Image, region_top, region_bottom, region_right,
                     min_height, min_conf) -> str:
    """OCR one slide image and return only the text inside the main content
    region (the overlay strips — ribbon, side panel, taskbar — are excluded)."""
    global _ocr_reader
    if _ocr_reader is None:
        import easyocr
        _ocr_reader = easyocr.Reader(["ch_sim", "en"])
    W, H = img.size
    arr = np.array(img)
    lines = []
    for box, txt, conf in _ocr_reader.readtext(arr):
        xs = [p[0] for p in box]
        ys = [p[1] for p in box]
        cx = (max(xs) + min(xs)) / 2 / W
        cy = (max(ys) + min(ys)) / 2 / H
        bh = (max(ys) - min(ys)) / H
        if cx > region_right:
            continue
        if cy < region_top or cy > region_bottom:
            continue
        if bh < min_height:
            continue
        if conf < min_conf:
            continue
        lines.append(txt)
    return "\n".join(lines)


def main():
    parser = argparse.ArgumentParser(description="Deduplicate an existing pptx")
    parser.add_argument("input_pptx", type=Path)
    parser.add_argument("output_pptx", type=Path)
    parser.add_argument("--image-similarity", type=float, default=0.95,
                        help="Drop a slide when its frame-image hash similarity to the "
                             "previous kept slide is >= this (0 disables)")
    parser.add_argument("--crop-top", type=float, default=0.0)
    parser.add_argument("--crop-bottom", type=float, default=0.0)
    parser.add_argument("--crop-left", type=float, default=0.0)
    parser.add_argument("--crop-right", type=float, default=0.0)
    parser.add_argument("--text-similarity", type=float, default=0.5,
                        help="Drop a slide when most of its OCR text is the same as the "
                             "previous kept slide's (similarity >= this, e.g. 0.50); "
                             "0 disables the OCR pass")
    parser.add_argument("--region-top", type=float, default=0.15,
                        help="Ignore OCR text above this fraction of height (ribbon/title bar)")
    parser.add_argument("--region-bottom", type=float, default=0.92,
                        help="Ignore OCR text below this fraction of height (taskbar)")
    parser.add_argument("--region-right", type=float, default=0.90,
                        help="Ignore OCR text right of this fraction of width (side panel)")
    parser.add_argument("--min-text-height", type=float, default=0.01,
                        help="Ignore OCR lines smaller than this fraction of height")
    parser.add_argument("--min-conf", type=float, default=0.3,
                        help="Ignore OCR lines with confidence below this")
    args = parser.parse_args()

    crop = (args.crop_top, args.crop_bottom, args.crop_left, args.crop_right)
    crop = crop if any(crop) else None

    src = Presentation(str(args.input_pptx))
    total = len(src.slides)

    slides = []
    for s in src.slides:
        img = None
        for sh in s.shapes:
            if sh.shape_type == 13:  # PICTURE
                img = Image.open(BytesIO(sh.image.blob)).convert("RGB")
                break
        notes = s.notes_slide.notes_text_frame.text if s.has_notes_slide else ""
        slides.append({"image": img, "notes": notes})

    keep_idx, last_hash, last_text = [], None, None
    for i, sl in enumerate(slides):
        drop = False
        if args.image_similarity > 0 and sl["image"] is not None and last_hash is not None:
            h = image_hash_pil(sl["image"], crop=crop)
            if hash_similarity(h, last_hash) >= args.image_similarity:
                drop = True
        if not drop and args.text_similarity > 0 and sl["image"] is not None:
            txt = ocr_content_text(sl["image"], args.region_top, args.region_bottom,
                                   args.region_right, args.min_text_height, args.min_conf)
            if last_text is not None and text_similarity(txt, last_text) >= args.text_similarity:
                drop = True
            if not drop:
                last_text = txt
        if drop:
            continue
        keep_idx.append(i)
        if args.image_similarity > 0 and sl["image"] is not None:
            last_hash = image_hash_pil(sl["image"], crop=crop)

    out = Presentation()
    out.slide_width = src.slide_width
    out.slide_height = src.slide_height
    blank = out.slide_layouts[6]
    for i in keep_idx:
        s = src.slides[i]
        slide = out.slides.add_slide(blank)
        for sh in s.shapes:
            if sh.shape_type == 13:
                img = sh.image
                slide.shapes.add_picture(BytesIO(img.blob), sh.left, sh.top,
                                         sh.width, sh.height)
                break
        if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip():
            slide.notes_slide.notes_text_frame.text = s.notes_slide.notes_text_frame.text

    out.save(str(args.output_pptx))
    print(f"{total} slides -> {len(keep_idx)} slides "
          f"(image-similarity={args.image_similarity}, crop={crop}, "
          f"text-similarity={args.text_similarity}, "
          f"region=(top={args.region_top}, bottom={args.region_bottom}, "
          f"right={args.region_right}))")
    print(f"Saved {args.output_pptx}")


if __name__ == "__main__":
    main()
