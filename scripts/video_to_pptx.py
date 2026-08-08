#!/usr/bin/env python3
"""
video_to_pptx.py
Convert a video lecture/courseware into a .pptx deck.

Steps:
  1. Extract frames at a regular interval with ffmpeg.
  2. Capture gate: skip frames whose hash similarity to the previous captured
     frame is above --sample-similarity (default 0.5).
  3. Region dedup: drop adjacent frames where > --region-ratio of the
     grid*grid regions (default 8x8 = 64) are "the same" as the previous kept
     frame (region hash similarity > --region-similarity).
  4. Deduplicate visually similar frames via global perceptual hashing
     (--similarity), then optionally merge frames with identical OCR text.
  5. OCR each kept frame (PaddleOCR preferred, fallback to easyocr).
  6. Build one slide per unique frame, embedding the frame image and
     appending recognized text to the slide notes.

Usage:
  python scripts/video_to_pptx.py input.mp4 output.pptx \
      --interval 2 --similarity 0.92 --ocr easyocr \
      --crop-top 0.05 --crop-bottom 0.06 --text-similarity 0.9
"""
import argparse
import difflib
import hashlib
import re
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path

import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


def run(cmd):
    print(f"$ {' '.join(cmd)}")
    subprocess.run(cmd, check=True)


def extract_frames(video_path: Path, out_dir: Path, interval_sec: float):
    out_dir.mkdir(parents=True, exist_ok=True)
    fps = f"1/{interval_sec}"
    cmd = [
        "ffmpeg", "-y", "-i", str(video_path),
        "-vf", f"fps={fps},scale='min(iw,1920)':-1",
        "-q:v", "2",
        str(out_dir / "frame_%06d.jpg"),
    ]
    run(cmd)
    return sorted(out_dir.glob("frame_*.jpg"))


def pixel_hash(img: Image.Image, size: int = 16, crop=None) -> str:
    """Perceptual hash of an already-open PIL image. `crop` is (top, bottom,
    left, right) fractions of the image to ignore before hashing — used to
    remove overlay regions such as title bars, taskbars or player controls
    from the comparison."""
    g = img.convert("L")
    if crop:
        w, h = g.size
        top, bottom, left, right = crop
        box = (
            int(w * left), int(h * top),
            int(w * (1 - right)), int(h * (1 - bottom)),
        )
        g = g.crop(box)
    g = g.resize((size, size), Image.Resampling.LANCZOS)
    arr = np.array(g, dtype=np.float32)
    bits = (arr > arr.mean()).astype(np.uint8)
    return hashlib.sha1(bits.tobytes()).hexdigest()


def image_hash(path: Path, size: int = 16, crop=None) -> str:
    return pixel_hash(Image.open(path), size=size, crop=crop)


def hash_similarity(a: str, b: str) -> float:
    bin_a = bin(int(a, 16))[2:].zfill(160)
    bin_b = bin(int(b, 16))[2:].zfill(160)
    dist = sum(x != y for x, y in zip(bin_a, bin_b))
    return 1 - dist / 160.0


def deduplicate(frame_paths, threshold: float, crop=None):
    kept, hashes = [], []
    for p in frame_paths:
        h = image_hash(p, crop=crop)
        if any(hash_similarity(h, e) >= threshold for e in hashes):
            continue
        kept.append(p)
        hashes.append(h)
    return kept


def sample_gate(frame_paths, threshold: float, crop=None):
    """Capture rule: sample every --interval, but skip a frame whose global hash
    similarity to the previous captured frame is above `threshold` (0.5 = "same
    screen, don't capture"). Set threshold <= 0 to disable."""
    if threshold <= 0:
        return frame_paths
    kept, last = [], None
    for p in frame_paths:
        h = pixel_hash(Image.open(p), crop=crop)
        if last is not None and hash_similarity(h, last) > threshold:
            continue
        kept.append(p)
        last = h
    return kept


def region_signature(img: Image.Image, grid: int, crop=None) -> list:
    """Per-region perceptual hashes for a grid x grid division of the frame,
    returned row-major (grid*grid entries). `crop` strips overlay strips
    (title bar / taskbar / side panel) before the grid is laid out."""
    g = img.convert("L")
    if crop:
        w, h = g.size
        top, bottom, left, right = crop
        g = g.crop((int(w * left), int(h * top),
                    int(w * (1 - right)), int(h * (1 - bottom))))
    w, h = g.size
    cw, ch = w // grid, h // grid
    sig = []
    for r in range(grid):
        for c in range(grid):
            box = (c * cw, r * ch,
                   w if c == grid - 1 else (c + 1) * cw,
                   h if r == grid - 1 else (r + 1) * ch)
            sig.append(pixel_hash(g.crop(box)))
    return sig


def region_dedup(frame_paths, grid: int, region_sim: float, region_ratio: float, crop=None):
    """Adjacent-frame dedup by region: divide each frame into grid*grid regions
    and compare against the previous kept frame. A region counts as "the same"
    when its hash similarity is above `region_sim`; when more than `region_ratio`
    (fraction) of the regions are the same, the frame is a duplicate and dropped.
    Set region_ratio <= 0 to disable."""
    if region_ratio <= 0 or grid <= 0:
        return frame_paths
    kept, sigs = [], []
    for p in frame_paths:
        sig = region_signature(Image.open(p), grid, crop=crop)
        if sigs:
            same = sum(1 for a, b in zip(sig, sigs[-1])
                       if hash_similarity(a, b) > region_sim)
            if same / (grid * grid) > region_ratio:
                continue
        kept.append(p)
        sigs.append(sig)
    return kept


def text_similarity(a: str, b: str) -> float:
    """0-1 similarity of two OCR text blocks, whitespace-insensitive."""
    a = re.sub(r"\s+", "", a)
    b = re.sub(r"\s+", "", b)
    if not a or not b:
        return 0.0
    return difflib.SequenceMatcher(None, a, b).ratio()


def merge_by_text(frame_paths, texts, threshold: float):
    """Drop a frame when its main content (OCR text) is nearly identical to the
    previously kept frame's. Implements "same main content => duplicate"."""
    if threshold <= 0 or len(frame_paths) < 2:
        return frame_paths, texts
    kept_frames, kept_texts, last = [], [], None
    for f, t in zip(frame_paths, texts):
        if last is not None and text_similarity(t, last) >= threshold:
            continue
        kept_frames.append(f)
        kept_texts.append(t)
        last = t
    return kept_frames, kept_texts


def paddle_major() -> int:
    """Return installed paddleocr major version (default 2 if undetectable)."""
    try:
        from importlib.metadata import version
        return int(version("paddleocr").split(".")[0])
    except Exception:
        return 2


def paddle_ocr():
    """Build a cached PaddleOCR instance matching the installed API version."""
    from paddleocr import PaddleOCR
    ocr = getattr(paddle_ocr, "_inst", None)
    if ocr is None:
        if paddle_major() >= 3:
            # PaddleOCR 3.x: ocr()/use_angle_cls are deprecated; use predict().
            # Disable orientation/doc-unwarp preprocessing to speed up slide OCR.
            try:
                ocr = PaddleOCR(
                    lang="ch",
                    use_doc_orientation_classify=False,
                    use_doc_unwarping=False,
                    use_textline_orientation=False,
                )
            except (TypeError, ValueError):
                ocr = PaddleOCR(lang="ch")
        else:
            # PaddleOCR 2.x legacy API.
            ocr = PaddleOCR(use_angle_cls=True, lang="ch", show_log=False)
        paddle_ocr._inst = ocr
    return ocr


def _filter_ocr_lines(path: Path, rows, min_height=0.0, text_top=0.0, text_bottom=0.0) -> list:
    """Keep only the frame's main-content text:
    - drop lines smaller than min_height fraction of frame height
    - drop lines whose box center sits in the top/bottom overlay strips
      (title bars, ribbon, taskbar, watermarks), defined by text_top/text_bottom
      fractions of height."""
    img_h = Image.open(path).height
    out = []
    for t, box in rows:
        ys = [p[1] for p in box]
        cy = (max(ys) + min(ys)) / 2 / img_h
        bh = (max(ys) - min(ys)) / img_h
        if text_top and cy < text_top:
            continue
        if text_bottom and cy > 1 - text_bottom:
            continue
        if min_height and bh < min_height:
            continue
        out.append(t)
    return out


def ocr_text_paddle(path: Path, min_height: float = 0.0,
                    text_top: float = 0.0, text_bottom: float = 0.0) -> str:
    ocr = paddle_ocr()
    if paddle_major() >= 3:
        texts = []
        for res in ocr.predict(str(path)):
            rec = getattr(res, "rec_texts", None) or res.get("rec_texts", []) or []
            boxes = getattr(res, "rec_boxes", None) or res.get("rec_boxes", None)
            if boxes is not None:
                rows = []
                for i, t in enumerate(rec):
                    b = np.asarray(boxes[i])
                    pts = b.reshape(-1, 2).tolist() if b.size else [[0, 0]]
                    rows.append((t, pts))
                texts.extend(_filter_ocr_lines(path, rows, min_height, text_top, text_bottom))
            else:
                texts.extend(rec)
        return "\n".join(texts)

    result = ocr.ocr(str(path), cls=True)
    if not result or not result[0]:
        return ""
    rows = [(line[1][0], line[0]) for line in result[0]]
    return "\n".join(_filter_ocr_lines(path, rows, min_height, text_top, text_bottom))


def ocr_text(path: Path, engine: str, min_height: float = 0.0,
             text_top: float = 0.0, text_bottom: float = 0.0) -> str:
    if engine == "paddleocr":
        return ocr_text_paddle(path, min_height, text_top, text_bottom)

    import easyocr
    reader = getattr(ocr_text, "_easyocr", None)
    if reader is None:
        reader = easyocr.Reader(["ch_sim", "en"])
        ocr_text._easyocr = reader
    rows = [(txt, box) for box, txt, _ in reader.readtext(str(path))]
    return "\n".join(_filter_ocr_lines(path, rows, min_height, text_top, text_bottom))


def build_pptx(frame_paths, texts, output: Path, size=(13.333, 7.5)):
    prs = Presentation()
    prs.slide_width = Inches(size[0])
    prs.slide_height = Inches(size[1])
    blank = prs.slide_layouts[6]

    for frame, text in zip(frame_paths, texts):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(frame), Inches(0), Inches(0),
            width=prs.slide_width, height=prs.slide_height
        )
        if text.strip():
            slide.notes_slide.notes_text_frame.text = text.strip()

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    print(f"Saved {output} with {len(frame_paths)} slides.")


def main():
    parser = argparse.ArgumentParser(description="Convert video courseware to PPTX")
    parser.add_argument("input_video", type=Path)
    parser.add_argument("output_pptx", type=Path)
    parser.add_argument("--interval", type=float, default=10.0,
                        help="Seconds between extracted frames")
    parser.add_argument("--similarity", type=float, default=0.92,
                        help="Perceptual-hash similarity threshold (0-1)")
    parser.add_argument("--sample-similarity", type=float, default=0.5,
                        help="Capture gate: skip a frame whose hash similarity to the "
                             "previous captured frame is above this (0 disables)")
    parser.add_argument("--region-grid", type=int, default=8,
                        help="Divide frames into grid x grid regions for adjacent-frame "
                             "dedup (8 -> 64 regions)")
    parser.add_argument("--region-similarity", type=float, default=0.5,
                        help="Per-region hash similarity above which a region counts "
                             "as \"the same\" between adjacent frames")
    parser.add_argument("--region-ratio", type=float, default=0.8,
                        help="Drop a frame when more than this fraction of its regions "
                             "match the previous kept frame (0 disables region dedup)")
    parser.add_argument("--ocr", choices=["paddleocr", "easyocr"], default="paddleocr")
    parser.add_argument("--keep-frames", action="store_true")
    parser.add_argument("--crop-top", type=float, default=0.0,
                        help="Fraction of height to ignore at top (title bars etc.)")
    parser.add_argument("--crop-bottom", type=float, default=0.0,
                        help="Fraction of height to ignore at bottom (taskbar/player bar)")
    parser.add_argument("--crop-left", type=float, default=0.0,
                        help="Fraction of width to ignore at left")
    parser.add_argument("--crop-right", type=float, default=0.0,
                        help="Fraction of width to ignore at right")
    parser.add_argument("--text-similarity", type=float, default=0.0,
                        help="Merge slides whose OCR text similarity >= this (0 disables)")
    parser.add_argument("--min-text-height", type=float, default=0.02,
                        help="Drop OCR lines smaller than this fraction of frame "
                             "height (removes UI/overlay text; 0 keeps all)")
    parser.add_argument("--text-top", type=float, default=0.0,
                        help="Fraction of height to ignore at top when keeping OCR "
                             "text (ribbon / title bar overlay removal)")
    parser.add_argument("--text-bottom", type=float, default=0.0,
                        help="Fraction of height to ignore at bottom when keeping OCR "
                             "text (taskbar / watermark overlay removal)")
    args = parser.parse_args()

    if not args.input_video.exists():
        sys.exit(f"Input not found: {args.input_video}")

    work = Path(tempfile.mkdtemp(prefix="v2pptx_"))
    try:
        frames_dir = work / "frames"
        print(f"Extracting frames every {args.interval}s...")
        frames = extract_frames(args.input_video, frames_dir, args.interval)

        crop = (args.crop_top, args.crop_bottom, args.crop_left, args.crop_right)
        crop = crop if any(crop) else None

        if args.sample_similarity > 0:
            before = len(frames)
            frames = sample_gate(frames, args.sample_similarity, crop=crop)
            print(f"Capture gate: {before} -> {len(frames)} "
                  f"(sim > {args.sample_similarity} to previous frame skipped)")

        if args.region_ratio > 0:
            before = len(frames)
            frames = region_dedup(frames, args.region_grid, args.region_similarity,
                                  args.region_ratio, crop=crop)
            print(f"Region dedup: {before} -> {len(frames)} "
                  f"({args.region_grid}x{args.region_grid} grid, "
                  f"region sim > {args.region_similarity}, ratio > {args.region_ratio})")

        print(f"Deduplicating {len(frames)} frames (threshold={args.similarity}, crop={crop})...")
        unique = deduplicate(frames, args.similarity, crop=crop)

        print(f"OCR {len(unique)} unique frames with {args.ocr} "
              f"(min-text-height={args.min_text_height}, "
              f"text-top={args.text_top}, text-bottom={args.text_bottom})...")
        texts = [ocr_text(f, args.ocr, args.min_text_height, args.text_top, args.text_bottom) for f in unique]

        if args.text_similarity > 0:
            before = len(unique)
            unique, texts = merge_by_text(unique, texts, args.text_similarity)
            print(f"Text dedup: {before} -> {len(unique)} (similarity>={args.text_similarity})")

        build_pptx(unique, texts, args.output_pptx)

        if args.keep_frames:
            keep = args.output_pptx.with_suffix(".frames")
            shutil.copytree(frames_dir, keep, dirs_exist_ok=True)
            print(f"Frames kept at {keep}")
    finally:
        if not args.keep_frames:
            shutil.rmtree(work, ignore_errors=True)


if __name__ == "__main__":
    main()
